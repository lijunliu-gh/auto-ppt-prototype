"""Python-native PPTX renderer using python-pptx.

Used when a brand template is provided. Falls back to JS renderer when no template.
"""

from __future__ import annotations

import io
import json
import logging
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .template_engine import TemplateConfig

logger = logging.getLogger("auto-ppt")

# Slide dimensions (widescreen 13.333 x 7.5 inches)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Default colors (hex strings without #)
_COLORS = {
    "primary": "0F766E",
    "secondary": "2563EB",
    "accent": "F59E0B",
    "danger": "DC2626",
    "bg_light": "F8FAFC",
    "bg_card": "F1F5F9",
    "border": "CBD5E1",
    "text_dark": "0F172A",
    "text_body": "1F2937",
    "text_muted": "475569",
    "text_subtle": "64748B",
    "text_label": "334155",
    "header_fill": "E2E8F0",
    "closing_bg": "0F172A",
    "closing_text": "F8FAFC",
    "closing_sub": "CBD5E1",
    "green_light": "ECFDF5",
    "blue_light": "EFF6FF",
}
_CHART_COLORS = ["0F766E", "2563EB", "F59E0B", "DC2626"]


def _rgb(hex_str: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _get_colors(template: Optional[TemplateConfig]) -> Dict[str, str]:
    """Get color palette, merging template theme colors if available."""
    colors = dict(_COLORS)
    if template and template.theme_colors:
        tc = template.theme_colors
        colors["primary"] = tc.primary
        colors["secondary"] = tc.secondary
        colors["accent"] = tc.accent
        colors["text_dark"] = tc.text_dark
        colors["text_body"] = tc.text_dark  # body inherits from text_dark
        colors["closing_text"] = tc.text_light
        colors["bg_light"] = tc.background
    return colors


def render_deck_with_template(
    deck: Dict[str, Any],
    output_json_path: str | Path,
    output_pptx_path: str | Path,
    template: TemplateConfig,
    base_dir: str | Path | None = None,
) -> None:
    """Render a deck JSON to PPTX using a brand template.

    Args:
        deck: Validated deck JSON dict.
        output_json_path: Where to save the deck JSON.
        output_pptx_path: Where to save the output PPTX.
        template: Parsed template configuration.
        base_dir: Base directory for resolving relative image paths.
    """
    json_path = Path(output_json_path)
    pptx_path = Path(output_pptx_path)

    json_path.parent.mkdir(parents=True, exist_ok=True)
    pptx_path.parent.mkdir(parents=True, exist_ok=True)

    # Save the deck JSON
    json_path.write_text(json.dumps(deck, indent=2, ensure_ascii=False), encoding="utf-8")

    # Create presentation from template
    prs = Presentation(str(template.source_path))
    colors = _get_colors(template)
    fonts = (template.font_heading, template.font_body)

    # Remove any existing slides from the template (we build fresh)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get("r:id")
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # Render each slide
    resolved_base = Path(base_dir).resolve() if base_dir else json_path.parent.resolve()
    for slide_def in deck.get("slides", []):
        _render_slide(prs, deck, slide_def, template, colors, fonts, resolved_base)

    # Set presentation metadata
    prs.core_properties.author = "auto-ppt-engine"
    prs.core_properties.title = deck.get("deckTitle", "Untitled")
    prs.core_properties.subject = deck.get("scenario", "Auto-generated presentation")

    prs.save(str(pptx_path))
    logger.info("Rendered PPTX (template): %s -> %s", json_path.name, pptx_path.name)


def _get_layout(prs: Presentation, template: TemplateConfig, layout_name: str):
    """Get the template slide layout for a given schema layout name."""
    idx = template.layout_mapping.get(layout_name)
    if idx is not None and idx < len(prs.slide_layouts):
        return prs.slide_layouts[idx]
    # Fallback to first layout
    return prs.slide_layouts[0]


def _render_slide(
    prs: Presentation,
    deck: Dict[str, Any],
    slide_def: Dict[str, Any],
    template: TemplateConfig,
    colors: Dict[str, str],
    fonts: tuple,
    base_dir: Optional[Path] = None,
) -> None:
    """Render a single slide based on its layout type."""
    layout_name = slide_def.get("layout", "bullet")
    layout = _get_layout(prs, template, layout_name)
    slide = prs.slides.add_slide(layout)

    renderer = _LAYOUT_RENDERERS.get(layout_name, _render_bullet)
    renderer(slide, deck, slide_def, colors, fonts, base_dir=base_dir)

    # Speaker notes
    _add_notes(slide, deck, slide_def)

    # Footer page number (skip for title/closing)
    if layout_name not in ("title", "closing"):
        _add_footer(slide, slide_def.get("page", 1), colors, fonts)


def _fill_placeholder(slide, idx: int, text: str) -> bool:
    """Try to fill a placeholder by index. Returns True if successful."""
    try:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == idx:
                ph.text = text
                return True
    except Exception:
        pass
    return False


def _add_textbox(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str = "Aptos",
    font_size: int = 18,
    bold: bool = False,
    italic: bool = False,
    color: str = "1F2937",
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> None:
    """Add a text box to a slide."""
    if not text:
        return
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = _rgb(color)
    p.alignment = alignment

    # Set vertical alignment
    from pptx.oxml.ns import qn
    txBody = txBox.text_frame._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        anchor_map = {
            MSO_ANCHOR.TOP: "t",
            MSO_ANCHOR.MIDDLE: "ctr",
            MSO_ANCHOR.BOTTOM: "b",
        }
        bodyPr.set("anchor", anchor_map.get(valign, "t"))


def _add_bullet_list(
    slide,
    bullets: List[str],
    left: float,
    top: float,
    width: float,
    height: float,
    font_name: str = "Aptos",
    font_size: int = 18,
    color: str = "1F2937",
) -> None:
    """Add a bulleted list to a slide."""
    items = [b for b in (bullets or []) if b]
    if not items:
        return
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.color.rgb = _rgb(color)
        p.space_after = Pt(10)
        # Add bullet
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022"})
        # Remove existing buNone/buChar if any
        for child in list(pPr):
            if child.tag.endswith(("buChar", "buNone", "buAutoNum")):
                pPr.remove(child)
        pPr.append(buChar)


def _add_slide_header(
    slide, deck: Dict, slide_def: Dict, colors: Dict, fonts: tuple
) -> None:
    """Add standard content slide header with title, objective, deck title, divider."""
    _add_textbox(
        slide, slide_def.get("title", ""),
        0.6, 0.35, 8.3, 0.5,
        font_name=fonts[0], font_size=24, bold=True, color=colors["text_dark"],
    )
    if slide_def.get("objective"):
        _add_textbox(
            slide, slide_def["objective"],
            0.6, 0.88, 8.5, 0.35,
            font_name=fonts[1], font_size=10, italic=True, color=colors["text_muted"],
        )
    _add_textbox(
        slide, deck.get("deckTitle", ""),
        10.3, 0.36, 2.3, 0.25,
        font_name=fonts[1], font_size=9, color=colors["text_subtle"],
        alignment=PP_ALIGN.RIGHT,
    )
    # Divider line
    from pptx.util import Inches as In
    slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        In(0.6), In(1.28), In(12.7), In(1.28)
    ).line.color.rgb = _rgb(colors["border"])


def _add_label(
    slide, text: str, x: float, y: float, w: float, colors: Dict, fonts: tuple
) -> None:
    _add_textbox(
        slide, text, x, y, w, 0.25,
        font_name=fonts[1], font_size=11, bold=True, color=colors["text_label"],
    )


def _add_footer(slide, page: int, colors: Dict, fonts: tuple) -> None:
    _add_textbox(
        slide, str(page),
        12.3, 7.02, 0.35, 0.2,
        font_name=fonts[1], font_size=9, color=colors["text_subtle"],
        alignment=PP_ALIGN.RIGHT,
    )


def _add_notes(slide, deck: Dict, slide_def: Dict) -> None:
    """Add speaker notes and source citations."""
    lines = []
    if deck.get("needsSpeakerNotes") and slide_def.get("speakerNotes"):
        lines.extend(slide_def["speakerNotes"])

    if deck.get("sourceDisplayMode") == "notes" and slide_def.get("sources"):
        if lines:
            lines.append("")
        lines.append("Sources:")
        for i, src in enumerate(slide_def["sources"]):
            parts = [f"[{i+1}] {src.get('label', 'Source')}"]
            if src.get("location"):
                parts.append(src["location"])
            if src.get("citation"):
                parts.append(f"Citation: {src['citation']}")
            if src.get("notes"):
                parts.append(f"Note: {src['notes']}")
            lines.append(" | ".join(parts))

    if lines:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "\n".join(lines)


def _add_rect(
    slide, x: float, y: float, w: float, h: float,
    fill_color: Optional[str] = None,
    line_color: Optional[str] = None,
    radius: float = 0,
) -> None:
    """Add a rectangle shape to the slide."""
    from pptx.enum.shapes import MSO_SHAPE
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius > 0 else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill_color)
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = _rgb(line_color)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()


# ── Image Insertion ───────────────────────────────────────────────

def _insert_visuals_on_slide(
    slide,
    visuals: List,
    base_dir: Optional[Path],
    colors: Dict[str, str],
    fonts: tuple,
) -> bool:
    """Process visuals list: insert images, render placeholders, or show text suggestions.

    Returns True if at least one image was inserted.
    """
    from .image_handler import partition_visuals, resolve_image, get_position, classify_visual

    descriptions, images, placeholders = partition_visuals(visuals)
    inserted_any = False

    # Insert actual images
    for img_visual in images:
        if base_dir is None:
            continue
        result = resolve_image(img_visual, base_dir)
        if result:
            data, ext = result
            left, top, width, height = get_position(img_visual)
            _insert_image_bytes(slide, data, left, top, width, height)
            inserted_any = True
        else:
            # Failed to load → treat as description
            alt = img_visual.get("alt") or img_visual.get("path") or img_visual.get("url") or "Image"
            descriptions.append({"kind": "description", "text": f"[Image: {alt}]"})

    # Render placeholders as labeled boxes
    for ph in placeholders:
        left, top, width, height = get_position(ph)
        _add_rect(slide, left, top, width, height, fill_color=colors["bg_card"], line_color=colors["border"], radius=0.05)
        prompt_text = ph.get("prompt", "Image placeholder")
        _add_textbox(
            slide, f"🖼 {prompt_text}",
            left + 0.15, top + 0.15, width - 0.3, height - 0.3,
            font_name=fonts[1], font_size=11, italic=True, color=colors["text_muted"],
        )

    # Show remaining text descriptions as suggestion box (only if no images inserted)
    if descriptions and not inserted_any:
        desc_texts = [d.get("text", "") for d in descriptions if d.get("text")]
        if desc_texts:
            _add_rect(slide, 8.95, 1.65, 3.35, 3.2, fill_color=colors["bg_card"], line_color=colors["border"])
            _add_label(slide, "Visual Suggestions", 9.2, 1.92, 2.7, colors, fonts)
            _add_bullet_list(slide, desc_texts, 9.08, 2.25, 2.9, 2.2, fonts[1], 12, colors["text_body"])

    return inserted_any


def _insert_image_bytes(
    slide, data: bytes, left: float, top: float, width: float, height: float
) -> None:
    """Insert image bytes onto a slide."""
    image_stream = io.BytesIO(data)
    slide.shapes.add_picture(
        image_stream, Inches(left), Inches(top), Inches(width), Inches(height)
    )


# ── Layout Renderers ──────────────────────────────────────────────

def _render_title(slide, deck: Dict, sd: Dict, colors: Dict, fonts: tuple, **kwargs) -> None:
    """Title slide: large title, subtitle, metadata."""
    # Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(colors["bg_light"])

    # Accent bar
    _add_rect(slide, 0.6, 0.9, 0.18, 4.8, fill_color=colors["primary"])

    # Title
    _add_textbox(
        slide, sd.get("title", ""),
        1.0, 1.1, 8.8, 1.1,
        font_name=fonts[0], font_size=28, bold=True, color=colors["text_dark"],
    )
    # Subtitle
    subtitle = sd.get("subtitle") or deck.get("scenario") or deck.get("audience") or ""
    _add_textbox(
        slide, subtitle,
        1.02, 2.35, 7.4, 0.5,
        font_name=fonts[1], font_size=16, color=colors["text_label"],
    )
    # Meta
    meta = []
    if deck.get("audience"):
        meta.append(f"Audience: {deck['audience']}")
    if deck.get("scenario"):
        meta.append(f"Scenario: {deck['scenario']}")
    if deck.get("tone"):
        meta.append(f"Tone: {deck['tone']}")
    if meta:
        _add_textbox(
            slide, "\n".join(meta),
            1.02, 4.9, 4.2, 1.2,
            font_name=fonts[1], font_size=12, color=colors["text_muted"],
        )


def _render_agenda(slide, deck: Dict, sd: Dict, colors: Dict, fonts: tuple, **kwargs) -> None:
    """Agenda / section / summary slide."""
    _add_slide_header(slide, deck, sd, colors, fonts)
    items = [f"{i+1}. {b}" for i, b in enumerate(sd.get("bullets", []))]
    _add_bullet_list(slide, items, 0.95, 1.7, 7.8, 4.8, fonts[1], 20, colors["text_body"])


def _render_bullet(slide, deck: Dict, sd: Dict, colors: Dict, fonts: tuple, **kwargs) -> None:
    """Standard bullet slide with optional visuals — images, placeholders, or text."""
    _add_slide_header(slide, deck, sd, colors, fonts)
    _add_bullet_list(slide, sd.get("bullets", []), 0.9, 1.72, 7.5, 4.9, fonts[1], 20, colors["text_body"])

    visuals = sd.get("visuals", [])
    if not visuals:
        return

    base_dir = kwargs.get("base_dir")
    inserted = _insert_visuals_on_slide(slide, visuals, base_dir, colors, fonts)


def _render_two_column(slide, deck: Dict, sd: Dict, colors: Dict, fonts: tuple, **kwargs) -> None:
    _add_slide_header(slide, deck, sd, colors, fonts)
    _add_label(slide, "Left", 0.9, 1.55, 4.8, colors, fonts)
    _add_label(slide, "Right", 6.8, 1.55, 4.8, colors, fonts)
    # Divider
    slide.shapes.add_connector(
        1, Inches(6.45), Inches(1.75), Inches(6.45), Inches(6.4)
    ).line.color.rgb = _rgb(colors["border"])
    _add_bullet_list(slide, sd.get("left", []), 0.9, 1.9, 5.1, 4.8, fonts[1], 18, colors["text_body"])
    _add_bullet_list(slide, sd.get("right", []), 6.8, 1.9, 5.1, 4.8, fonts[1], 18, colors["text_body"])


def _render_comparison(slide, deck: Dict, sd: Dict, colors: Dict, fonts: tuple, **kwargs) -> None:
    _add_slide_header(slide, deck, sd, colors, fonts)
    _add_rect(slide, 0.85, 1.7, 5.45, 4.7, fill_color=colors["bg_light"], line_color=colors["border"], radius=0.05)
    _add_rect(slide, 6.95, 1.7, 5.45, 4.7, fill_color=colors["bg_light"], line_color=colors["border"], radius=0.05)
    _add_label(slide, "Option A", 1.1, 2.0, 4.5, colors, fonts)
    _add_label(slide, "Option B", 7.2, 2.0, 4.5, colors, fonts)
    _add_bullet_list(slide, sd.get("left", []), 1.0, 2.35, 4.9, 3.7, fonts[1], 16, colors["text_body"])
    _add_bullet_list(slide, sd.get("right", []), 7.1, 2.35, 4.9, 3.7, fonts[1], 16, colors["text_body"])


def _render_timeline(slide, deck: Dict, sd: Dict, colors: Dict, fonts: tuple, **kwargs) -> None:
    _add_slide_header(slide, deck, sd, colors, fonts)
    milestones = (sd.get("bullets") or [])[:5]
    if not milestones:
        return
    start_x = 1.1
    total_w = 10.8
    gap = total_w / (len(milestones) - 1) if len(milestones) > 1 else 0

    # Horizontal line
    slide.shapes.add_connector(
        1, Inches(start_x), Inches(3.65), Inches(start_x + total_w), Inches(3.65)
    ).line.color.rgb = _rgb("94A3B8")

    from pptx.enum.shapes import MSO_SHAPE
    for i, item in enumerate(milestones):
        cx = start_x + gap * i
        # Circle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(cx - 0.16), Inches(3.48), Inches(0.32), Inches(0.32),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(colors["primary"])
        shape.line.fill.background()

        # Phase label
        _add_textbox(
            slide, f"Phase {i+1}",
            cx - 0.45, 2.82, 0.9, 0.25,
            font_name=fonts[1], font_size=10, bold=True, color=colors["text_dark"],
            alignment=PP_ALIGN.CENTER,
        )
        # Description
        _add_textbox(
            slide, item,
            cx - 0.85, 3.95, 1.7, 1.0,
            font_name=fonts[1], font_size=11, color=colors["text_label"],
            alignment=PP_ALIGN.CENTER,
        )


def _render_process(slide, deck: Dict, sd: Dict, colors: Dict, fonts: tuple, **kwargs) -> None:
    _add_slide_header(slide, deck, sd, colors, fonts)
    steps = (sd.get("bullets") or [])[:4]
    if not steps:
        return
    box_w = 2.7

    for i, item in enumerate(steps):
        x = 0.9 + i * 3.0
        fill = colors["green_light"] if i % 2 == 0 else colors["blue_light"]
        _add_rect(slide, x, 2.2, box_w, 1.6, fill_color=fill, line_color=colors["border"], radius=0.05)
        # Step number
        _add_textbox(
            slide, str(i + 1),
            x + 0.15, 2.35, 0.3, 0.3,
            font_name=fonts[1], font_size=12, bold=True, color=colors["primary"],
        )
        # Step text
        _add_textbox(
            slide, item,
            x + 0.25, 2.68, 2.15, 0.85,
            font_name=fonts[1], font_size=15, bold=True, color=colors["text_dark"],
            alignment=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
        )


def _render_table(slide, deck: Dict, sd: Dict, colors: Dict, fonts: tuple, **kwargs) -> None:
    _add_slide_header(slide, deck, sd, colors, fonts)
    table_data = sd.get("table") or {}
    columns = table_data.get("columns", [])
    rows = table_data.get("rows", [])

    if not columns:
        _add_textbox(
            slide, "No table data supplied.",
            0.95, 2.0, 4.5, 0.35,
            font_name=fonts[1], font_size=18, color=colors["danger"],
        )
        return

    row_count = len(rows) + 1  # +1 for header
    col_count = len(columns)
    tbl = slide.shapes.add_table(
        row_count, col_count,
        Inches(0.9), Inches(1.8), Inches(11.6), Inches(min(4.8, row_count * 0.42)),
    ).table

    # Header row
    for ci, col_name in enumerate(columns):
        cell = tbl.cell(0, ci)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(colors["header_fill"])
        for p in cell.text_frame.paragraphs:
            p.font.name = fonts[1]
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = _rgb(colors["text_body"])

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            if ci >= col_count:
                break
            cell = tbl.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.name = fonts[1]
                p.font.size = Pt(13)
                p.font.color.rgb = _rgb(colors["text_body"])


def _render_chart(slide, deck: Dict, sd: Dict, colors: Dict, fonts: tuple, **kwargs) -> None:
    _add_slide_header(slide, deck, sd, colors, fonts)
    chart_data_def = sd.get("chart") or {}
    categories = chart_data_def.get("categories", [])
    series_list = chart_data_def.get("series", [])

    if not categories or not series_list:
        # Placeholder box for empty chart
        _add_rect(slide, 1.0, 1.9, 11.2, 3.8, fill_color=colors["bg_light"], line_color=colors["border"], radius=0.06)
        _add_textbox(
            slide, chart_data_def.get("title") or "Chart Placeholder",
            1.3, 2.25, 4.0, 0.35,
            font_name=fonts[0], font_size=18, bold=True, color=colors["text_dark"],
        )
        _add_bullet_list(
            slide,
            sd.get("visuals") or ["Provide categories and series data to render a real chart."],
            1.3, 2.8, 9.4, 2.2, fonts[1], 16, colors["text_body"],
        )
        return

    chart_type_str = (chart_data_def.get("type") or "bar").lower()
    xl_chart_type = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "area": XL_CHART_TYPE.AREA,
    }.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series_list:
        data = s.get("data", [])
        # Ensure data length matches categories
        while len(data) < len(categories):
            data.append(0)
        chart_data.add_series(s.get("name", "Series"), data[:len(categories)])

    chart_frame = slide.shapes.add_chart(
        xl_chart_type, Inches(0.95), Inches(1.8), Inches(8.0), Inches(4.5),
        chart_data,
    )
    chart = chart_frame.chart
    chart.has_legend = True

    # Apply chart colors
    plot = chart.plots[0]
    for si, series_obj in enumerate(plot.series):
        color_hex = _CHART_COLORS[si % len(_CHART_COLORS)]
        series_obj.format.fill.solid()
        series_obj.format.fill.fore_color.rgb = _rgb(color_hex)

    # Takeaways box
    bullets = sd.get("bullets", [])
    if bullets:
        _add_rect(slide, 9.35, 1.85, 2.9, 3.6, fill_color=colors["bg_light"], line_color=colors["border"], radius=0.05)
        _add_label(slide, "Takeaways", 9.6, 2.1, 2.2, colors, fonts)
        _add_bullet_list(slide, bullets, 9.45, 2.45, 2.45, 2.7, fonts[1], 12, colors["text_body"])


def _render_quote(slide, deck: Dict, sd: Dict, colors: Dict, fonts: tuple, **kwargs) -> None:
    _add_slide_header(slide, deck, sd, colors, fonts)
    _add_rect(slide, 1.0, 1.95, 11.0, 3.2, fill_color=colors["bg_light"], line_color=colors["border"], radius=0.08)
    quote_text = (sd.get("bullets") or [sd.get("subtitle") or ""])[0] or ""
    _add_textbox(
        slide, f'"{quote_text}"',
        1.5, 2.45, 10.0, 1.2,
        font_name=fonts[1], font_size=24, italic=True, bold=True,
        color=colors["text_dark"], alignment=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
    )
    if sd.get("subtitle"):
        _add_textbox(
            slide, sd["subtitle"],
            8.6, 4.55, 2.3, 0.25,
            font_name=fonts[1], font_size=11, color=colors["text_muted"],
            alignment=PP_ALIGN.RIGHT,
        )


def _render_closing(slide, deck: Dict, sd: Dict, colors: Dict, fonts: tuple, **kwargs) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(colors["closing_bg"])
    _add_textbox(
        slide, sd.get("title", ""),
        1.0, 2.0, 8.8, 0.9,
        font_name=fonts[0], font_size=30, bold=True, color=colors["closing_text"],
    )
    _add_textbox(
        slide, sd.get("subtitle") or "Questions and discussion",
        1.02, 3.1, 5.4, 0.45,
        font_name=fonts[1], font_size=16, color=colors["closing_sub"],
    )
    if sd.get("bullets"):
        _add_bullet_list(slide, sd["bullets"], 1.0, 4.1, 5.8, 1.8, fonts[1], 15, colors["closing_text"])


# Layout renderer dispatch table
_LAYOUT_RENDERERS = {
    "title": _render_title,
    "agenda": _render_agenda,
    "section": _render_agenda,
    "summary": _render_agenda,
    "bullet": _render_bullet,
    "two-column": _render_two_column,
    "comparison": _render_comparison,
    "timeline": _render_timeline,
    "process": _render_process,
    "table": _render_table,
    "chart": _render_chart,
    "quote": _render_quote,
    "closing": _render_closing,
}
