from __future__ import annotations

import json
import shutil
import subprocess
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from pptx import Presentation

EMU_PER_INCH = 914400


@dataclass
class Issue:
    code: str
    message: str
    severity: str = "warning"

    def to_dict(self) -> dict[str, str]:
        return {
            "code": self.code,
            "message": self.message,
            "severity": self.severity,
        }


def _emu_to_inches(value: int) -> float:
    return float(value) / EMU_PER_INCH


def _shape_bbox_in_inches(shape: Any) -> tuple[float, float, float, float] | None:
    try:
        x = _emu_to_inches(int(shape.left))
        y = _emu_to_inches(int(shape.top))
        w = _emu_to_inches(int(shape.width))
        h = _emu_to_inches(int(shape.height))
    except Exception:  # noqa: BLE001
        return None
    if w <= 0 or h <= 0:
        return None
    return (x, y, x + w, y + h)


def _boxes_overlap(a: tuple[float, float, float, float], b: tuple[float, float, float, float]) -> bool:
    return not (a[2] <= b[0] or b[2] <= a[0] or a[3] <= b[1] or b[3] <= a[1])


def _run_command(command: list[str]) -> tuple[bool, str]:
    try:
        subprocess.run(command, capture_output=True, check=True, text=True)
        return True, ""
    except FileNotFoundError:
        return False, f"Command not found: {command[0]}"
    except subprocess.CalledProcessError as exc:
        stderr = (exc.stderr or "").strip()
        stdout = (exc.stdout or "").strip()
        details = stderr if stderr else stdout
        return False, details or f"Command failed: {' '.join(command)}"


def _export_images(pptx_path: Path, images_dir: Path, dpi: int) -> tuple[list[str], list[str]]:
    notes: list[str] = []
    images_dir.mkdir(parents=True, exist_ok=True)

    soffice = shutil.which("soffice")
    pdftoppm = shutil.which("pdftoppm")
    if not soffice:
        notes.append("Skipped image export: `soffice` is not available in PATH.")
        return [], notes
    if not pdftoppm:
        notes.append("Skipped image export: `pdftoppm` is not available in PATH.")
        return [], notes

    pdf_path = images_dir / f"{pptx_path.stem}.pdf"
    ok, err = _run_command(
        [
            soffice,
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(images_dir),
            str(pptx_path),
        ]
    )
    if not ok:
        notes.append(f"PPTX to PDF conversion failed: {err}")
        return [], notes

    if not pdf_path.exists():
        notes.append("PPTX to PDF conversion finished but no PDF file was produced.")
        return [], notes

    prefix = images_dir / "slide"
    ok, err = _run_command([pdftoppm, "-jpeg", "-r", str(dpi), str(pdf_path), str(prefix)])
    if not ok:
        notes.append(f"PDF to image conversion failed: {err}")
        return [], notes

    images = sorted(images_dir.glob("slide-*.jpg"))
    if not images:
        notes.append("Image conversion ran but no slide images were generated.")
        return [], notes

    return [str(path) for path in images], notes


def analyze_visual_quality(pptx_path: Path, margin_in: float = 0.3) -> dict[str, Any]:
    prs = Presentation(str(pptx_path))
    slide_w = _emu_to_inches(int(prs.slide_width))
    slide_h = _emu_to_inches(int(prs.slide_height))

    slide_reports: list[dict[str, Any]] = []

    for index, slide in enumerate(prs.slides, start=1):
        issues: list[Issue] = []
        text_shape_count = 0
        media_shape_count = 0
        boxes: list[tuple[str, tuple[float, float, float, float]]] = []

        for shape in slide.shapes:
            bbox = _shape_bbox_in_inches(shape)
            name = getattr(shape, "name", "shape")

            if getattr(shape, "has_text_frame", False):
                raw_text = (shape.text or "").strip()
                if raw_text:
                    text_shape_count += 1

            if hasattr(shape, "image") or getattr(shape, "has_chart", False) or getattr(shape, "has_table", False):
                media_shape_count += 1

            if not bbox:
                continue

            boxes.append((name, bbox))
            left, top, right, bottom = bbox
            if left < margin_in or top < margin_in or (slide_w - right) < margin_in or (slide_h - bottom) < margin_in:
                issues.append(
                    Issue(
                        code="edge_margin",
                        message=(
                            f"Shape '{name}' is close to slide edge. "
                            f"Minimum margin target is {margin_in:.2f}in."
                        ),
                    )
                )

        for i in range(len(boxes)):
            name_a, box_a = boxes[i]
            for j in range(i + 1, len(boxes)):
                name_b, box_b = boxes[j]
                if _boxes_overlap(box_a, box_b):
                    issues.append(
                        Issue(
                            code="overlap_candidate",
                            message=f"Possible overlap between '{name_a}' and '{name_b}'.",
                        )
                    )

        if text_shape_count == 0 and media_shape_count == 0:
            issues.append(Issue(code="empty_slide", message="Slide appears to have no visible content.", severity="high"))
        if index not in (1, len(prs.slides)) and text_shape_count == 1 and media_shape_count == 0:
            issues.append(
                Issue(
                    code="text_only",
                    message="Slide is mostly text-only. Consider adding a visual element.",
                    severity="info",
                )
            )

        slide_reports.append(
            {
                "slide": index,
                "issues": [item.to_dict() for item in issues],
            }
        )

    return {
        "slideCount": len(prs.slides),
        "slides": slide_reports,
    }


def run_visual_qa(
    pptx_path: str | Path,
    output_dir: str | Path | None = None,
    dpi: int = 150,
    margin_in: float = 0.3,
) -> dict[str, Any]:
    target = Path(pptx_path).resolve()
    if not target.exists():
        raise RuntimeError(f"PPTX file not found: {target}")

    report_dir = Path(output_dir).resolve() if output_dir else (target.parent / f"{target.stem}-qa")
    report_dir.mkdir(parents=True, exist_ok=True)
    images_dir = report_dir / "images"

    images, notes = _export_images(target, images_dir, dpi)
    analysis = analyze_visual_quality(target, margin_in=margin_in)

    total_issues = sum(len(s["issues"]) for s in analysis["slides"])
    high_risk = [
        s["slide"]
        for s in analysis["slides"]
        if any(item["severity"] == "high" for item in s["issues"])
    ]

    report = {
        "pptxPath": str(target),
        "generatedAt": datetime.now(timezone.utc).isoformat(),
        "imagesDir": str(images_dir),
        "images": images,
        "notes": notes,
        "analysis": analysis,
        "summary": {
            "slideCount": analysis["slideCount"],
            "totalIssues": total_issues,
            "highRiskSlides": high_risk,
        },
    }

    report_path = report_dir / "visual-qa-report.json"
    report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    report["reportPath"] = str(report_path)
    return report